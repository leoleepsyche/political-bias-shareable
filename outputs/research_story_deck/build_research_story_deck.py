from __future__ import annotations

import csv
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT_DIR = ROOT / "outputs" / "research_story_deck"
FIG_DIR = ROOT / "outputs" / "figures_instruct_coefficients"
DETECTION_FIG_DIR = ROOT / "outputs" / "figures_instruct_detection"
STORY_FIG_DIR = ROOT / "outputs" / "figures_research_story"
OUT_PATH = OUT_DIR / "political_bias_research_story_skeleton.pptx"

WIDE = Inches(13.333)
HIGH = Inches(7.5)
ACCENT = RGBColor(31, 78, 121)
TEXT = RGBColor(25, 25, 25)
MUTED = RGBColor(92, 92, 92)
LIGHT = RGBColor(238, 242, 246)
LEFT_SAFE = RGBColor(0, 114, 178)
RIGHT_SAFE = RGBColor(230, 159, 0)


def add_textbox(slide, left, top, width, height, text, size=24, bold=False, color=TEXT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Aptos"
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, Inches(0.55), Inches(0.35), Inches(12.2), Inches(0.55), title, 28, True)
    line = slide.shapes.add_shape(1, Inches(0.55), Inches(1.03), Inches(1.25), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.color.rgb = ACCENT
    if subtitle:
        add_textbox(slide, Inches(0.55), Inches(1.12), Inches(12.1), Inches(0.35), subtitle, 12, False, MUTED)


def add_bullets(slide, left, top, width, height, bullets, size=17):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(8)
    return box


def add_numbered_cards(slide, cards):
    card_w = Inches(3.85)
    gap = Inches(0.33)
    left = Inches(0.7)
    for idx, (heading, body) in enumerate(cards):
        x = left + idx * (card_w + gap)
        shape = slide.shapes.add_shape(5, x, Inches(2.05), card_w, Inches(3.2))
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT
        shape.line.color.rgb = RGBColor(210, 218, 226)
        add_textbox(slide, x + Inches(0.22), Inches(2.25), card_w - Inches(0.44), Inches(0.35), f"{idx + 1}. {heading}", 18, True, ACCENT)
        add_textbox(slide, x + Inches(0.22), Inches(2.85), card_w - Inches(0.44), Inches(1.85), body, 15, False, TEXT)


def add_panel(slide, left, top, width, height, heading, body, fill=LIGHT, heading_color=ACCENT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = RGBColor(210, 218, 226)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.15), width - Inches(0.36), Inches(0.28), heading, 13, True, heading_color)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.52), width - Inches(0.36), height - Inches(0.65), body, 12, False, TEXT)
    return shape


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = ACCENT
    arrow.line.color.rgb = ACCENT
    return arrow


def add_line(slide, x1, y1, x2, y2, color=RGBColor(110, 110, 110), width=1.2):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(width)
    return line


def add_image(slide, path, left, top, width, height):
    if not path.exists():
        raise FileNotFoundError(path)
    slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_detection_comparison_table(slide, left, top, width, height):
    csv_path = DETECTION_FIG_DIR / "detection_metrics_instruct.csv"
    if not csv_path.exists():
        raise FileNotFoundError(csv_path)

    with csv_path.open(encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle))

    headers = ["Model", "Det.", "Best\nlayer", "Best\nAUC", "Agg.\nAUC"]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    col_widths = [Inches(1.1), Inches(0.55), Inches(0.72), Inches(0.72), Inches(0.72)]
    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = col_width

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(8.5)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, row in enumerate(rows, start=1):
        values = [
            row["model_label"],
            "L" if row["side"] == "left" else "R",
            row["best_layer"],
            f"{float(row['best_test_auc']):.3f}",
            f"{float(row['aggregation_test_auc']):.3f}",
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else RGBColor(248, 249, 250)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(8.7)
                    run.font.bold = col_idx == 1
                    if col_idx == 1 and value == "L":
                        run.font.color.rgb = LEFT_SAFE
                    elif col_idx == 1 and value == "R":
                        run.font.color.rgb = RIGHT_SAFE
                    else:
                        run.font.color.rgb = TEXT

    return table_shape


def load_story_detection_rows():
    csv_path = STORY_FIG_DIR / "story_detection_data.csv"
    if not csv_path.exists():
        raise FileNotFoundError(csv_path)
    with csv_path.open(encoding="utf-8", newline="") as handle:
        rows = list(csv.DictReader(handle))

    order = {"Qwen": 0, "Mistral": 1, "Llama-3": 2}
    side_order = {"left": 0, "right": 1}
    return sorted(rows, key=lambda row: (order[row["model"]], side_order[row["ideology"]]))


def add_story_detection_chart(slide, left, top, width, height):
    rows = load_story_detection_rows()
    by_model = {model: {} for model in ["Qwen", "Mistral", "Llama-3"]}
    for row in rows:
        by_model[row["model"]][row["ideology"]] = float(row["best_test_auc"])

    models = ["Qwen", "Mistral", "Llama-3"]
    chart_data = CategoryChartData()
    chart_data.categories = models
    chart_data.add_series("Left steering", [by_model[model]["left"] for model in models])
    chart_data.add_series("Right steering", [by_model[model]["right"] for model in models])

    chart_shape = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data)
    chart = chart_shape.chart
    chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.value_axis.minimum_scale = 0.84
    chart.value_axis.maximum_scale = 1.02
    chart.value_axis.major_unit = 0.02
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = LEFT_SAFE
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = RGBColor(255, 127, 14)
    for series in chart.series:
        series.has_data_labels = False
    return chart_shape


def add_story_detection_table(slide, left, top, width, height):
    rows = load_story_detection_rows()
    headers = ["Model", "Det.", "Best\nlayer", "Best\nAUC"]
    table_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, width, height)
    table = table_shape.table

    col_widths = [Inches(1.15), Inches(0.55), Inches(0.78), Inches(0.86)]
    for idx, col_width in enumerate(col_widths):
        table.columns[idx].width = col_width

    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.name = "Aptos"
                run.font.size = Pt(8.8)
                run.font.bold = True
                run.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, row in enumerate(rows, start=1):
        det = "L" if row["ideology"] == "left" else "R"
        values = [
            row["model"],
            det,
            row["best_layer"],
            f"{float(row['best_test_auc']):.3f}",
        ]
        for col_idx, value in enumerate(values):
            cell = table.cell(row_idx, col_idx)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(255, 255, 255) if row_idx % 2 else RGBColor(248, 249, 250)
            cell.margin_left = Inches(0.03)
            cell.margin_right = Inches(0.03)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER if col_idx else PP_ALIGN.LEFT
                for run in paragraph.runs:
                    run.font.name = "Aptos"
                    run.font.size = Pt(8.8)
                    run.font.bold = col_idx == 1
                    if col_idx == 1 and value == "L":
                        run.font.color.rgb = LEFT_SAFE
                    elif col_idx == 1 and value == "R":
                        run.font.color.rgb = RIGHT_SAFE
                    else:
                        run.font.color.rgb = TEXT
    return table_shape


def blank(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    return slide


def build():
    prs = Presentation()
    prs.slide_width = WIDE
    prs.slide_height = HIGH

    # 1
    s = blank(prs)
    add_textbox(s, Inches(0.75), Inches(1.35), Inches(11.9), Inches(0.9), "Political Stance Detection and Steering in Instruction-Tuned LLMs", 32, True)
    add_textbox(s, Inches(0.78), Inches(2.35), Inches(10.8), Inches(0.55), "A step-by-step research story from current experiments", 20, False, ACCENT)
    add_bullets(
        s,
        Inches(0.9),
        Inches(3.45),
        Inches(11.1),
        Inches(2.1),
        [
            "Core question: are political stance patterns detectable, steerable, and consistent across languages?",
            "Current focus: Qwen2.5-7B-Instruct, Mistral-7B-Instruct-v0.2, and Llama-3-8B-Instruct.",
            "Main emerging result: detectability and steerability are not the same property.",
        ],
        19,
    )

    # 2
    s = blank(prs)
    add_title(s, "Research Problem", "Move from observing bias to testing causal controllability.")
    add_bullets(
        s,
        Inches(0.75),
        Inches(1.85),
        Inches(11.7),
        Inches(4.6),
        [
            "Many evaluations ask whether LLMs show political bias in their answers.",
            "This project asks a stronger question: whether political stance behavior can be systematically moved by activation steering.",
            "The research object is not only the final answer, but the relationship between internal representation, intervention strength, language, and model family.",
            "The practical implication: a model can be highly separable in detection while still being weakly steerable in behavior.",
        ],
        20,
    )

    # 3
    s = blank(prs)
    add_title(s, "Three Research Questions")
    add_numbered_cards(
        s,
        [
            ("Detectability", "Can political stance patterns be separated reliably in the model's responses or representations?"),
            ("Steerability", "If we apply a left/right steering vector, do model answers move in the intended direction?"),
            ("Consistency", "Are these effects stable across models, languages, and steering coefficients?"),
        ],
    )

    # 4
    s = blank(prs)
    add_title(s, "Concrete Example: Training Classification Data", "The classifier is trained on paired left/right ideological responses, not on Political Compass outputs.")
    add_panel(
        s,
        Inches(0.65),
        Inches(1.45),
        Inches(3.45),
        Inches(1.85),
        "Shared instruction",
        "Would you say that Black people are treated less fairly than white people, white people are treated less fairly than Black people, or both are treated about equally?",
        RGBColor(246, 248, 250),
    )
    add_arrow(s, Inches(4.35), Inches(2.15), Inches(0.55), Inches(0.32))
    add_panel(
        s,
        Inches(5.2),
        Inches(1.35),
        Inches(3.25),
        Inches(1.45),
        "Left-labeled response",
        "Example response_text:\nBlack people are treated less fairly than white people in dealings with the police...\n\nFor left detector: y = 1\nFor right detector: y = 0",
        RGBColor(246, 248, 250),
    )
    add_panel(
        s,
        Inches(5.2),
        Inches(3.25),
        Inches(3.25),
        Inches(1.45),
        "Right-labeled response",
        "Example response_text:\nThe focus should be on the individual's actions, not their race...\n\nFor left detector: y = 0\nFor right detector: y = 1",
        RGBColor(246, 248, 250),
    )
    add_arrow(s, Inches(8.7), Inches(2.55), Inches(0.55), Inches(0.32))
    add_panel(
        s,
        Inches(9.55),
        Inches(1.75),
        Inches(3.05),
        Inches(2.25),
        "Training inputs",
        "Each paired instruction becomes two classifier inputs after prompt wrapping:\n\nWhat is your opinion on the following statement?\nStatement: {response_text}\nOpinion:",
        RGBColor(235, 243, 252),
    )
    add_textbox(
        s,
        Inches(0.8),
        Inches(5.35),
        Inches(11.8),
        Inches(0.65),
        "Canonical split: 600 train pairs, 120 validation pairs, and 120 test pairs. Each pair contains one left and one right response, so each detector sees 1200 train inputs.",
        16,
        False,
        MUTED,
    )
    add_textbox(
        s,
        Inches(0.8),
        Inches(6.28),
        Inches(11.8),
        Inches(0.28),
        "Source: IDEOINST dataset from Chen et al. (EMNLP 2024), processed locally as data/ideoinst_clean/ideology_840_opinion.csv; split/labels in pipeline/detect.py.",
        10,
        False,
        MUTED,
    )

    # 5
    s = blank(prs)
    add_title(s, "Probe Training: Learning a Political Direction", "This stage trains a supervised classifier on hidden activations.")
    add_textbox(s, Inches(0.25), Inches(1.28), Inches(1.0), Inches(0.22), "PROBE", 8, True, MUTED)

    add_panel(
        s,
        Inches(0.55),
        Inches(1.65),
        Inches(3.2),
        Inches(1.25),
        "Training input x",
        "Prompt-wrapped ideological response:\n\nWhat is your opinion on the following statement?\nStatement: {response_text}\nOpinion:",
        RGBColor(248, 249, 250),
    )
    add_panel(
        s,
        Inches(0.55),
        Inches(3.35),
        Inches(3.2),
        Inches(1.05),
        "Supervised label y",
        "Left detector: left response = 1, right response = 0\nRight detector: right response = 1, left response = 0",
        RGBColor(248, 249, 250),
    )

    add_arrow(s, Inches(3.95), Inches(2.45), Inches(0.62), Inches(0.35))
    llm = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.75), Inches(2.1), Inches(1.1), Inches(0.85))
    llm.fill.solid()
    llm.fill.fore_color.rgb = RGBColor(255, 255, 255)
    llm.line.color.rgb = RGBColor(120, 120, 120)
    add_textbox(s, Inches(5.02), Inches(2.36), Inches(0.55), Inches(0.22), "LLM", 15, False, TEXT)

    block_y = Inches(3.25)
    blocks = [
        ("Read-in", Inches(4.0), Inches(1.0)),
        ("Block 1", Inches(5.55), Inches(1.05)),
        ("Block 2", Inches(7.1), Inches(1.05)),
        ("...", Inches(8.55), Inches(0.55)),
        ("Block L", Inches(9.6), Inches(1.05)),
        ("Read-out", Inches(11.15), Inches(1.15)),
    ]
    add_line(s, Inches(5.3), Inches(2.95), Inches(5.3), block_y)
    prev_right = None
    for label, x, w in blocks:
        shape = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, block_y, w, Inches(0.62))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.color.rgb = RGBColor(145, 145, 145)
        add_textbox(s, x + Inches(0.08), block_y + Inches(0.2), w - Inches(0.16), Inches(0.18), label, 10, False, TEXT)
        if prev_right is not None:
            add_line(s, prev_right, block_y + Inches(0.31), x, block_y + Inches(0.31))
        prev_right = x + w

    for x, label in [(Inches(6.0), "A1"), (Inches(7.55), "A2"), (Inches(10.05), "AL")]:
        add_line(s, x, block_y + Inches(0.62), x + Inches(0.2), Inches(4.6))
        add_textbox(s, x + Inches(0.03), Inches(4.52), Inches(0.45), Inches(0.18), label, 9, False, MUTED)
        pred = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x - Inches(0.25), Inches(5.0), Inches(0.95), Inches(0.45))
        pred.fill.solid()
        pred.fill.fore_color.rgb = RGBColor(255, 255, 255)
        pred.line.color.rgb = RGBColor(150, 150, 150)
        add_textbox(s, x - Inches(0.08), Inches(5.13), Inches(0.62), Inches(0.14), "0 or 1?", 9, False, TEXT)

    add_panel(
        s,
        Inches(9.65),
        Inches(1.55),
        Inches(2.7),
        Inches(1.25),
        "Probe objective",
        "Train RFM / linear probe to predict y from layer activations.",
        RGBColor(235, 243, 252),
    )
    add_panel(
        s,
        Inches(9.65),
        Inches(5.1),
        Inches(2.7),
        Inches(0.95),
        "Model selection",
        "Choose l* by validation AUC; save that layer's direction for downstream steering.",
        RGBColor(235, 243, 252),
    )
    add_textbox(
        s,
        Inches(0.75),
        Inches(6.45),
        Inches(11.8),
        Inches(0.28),
        "Method reference: Beaglehole et al., Science 391(6787):787-792 (2026), doi:10.1126/science.aea6792; local implementation via neural_controllers and pipeline/detect.py.",
        10,
        False,
        MUTED,
    )

    # 6
    s = blank(prs)
    add_title(s, "Activation Steering: Mathematical Formulation", "The learned probe direction becomes an additive intervention in hidden-state space.")
    add_textbox(s, Inches(0.25), Inches(1.24), Inches(1.0), Inches(0.22), "STEER", 8, True, MUTED)

    formula_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.55), Inches(6.15), Inches(3.05))
    formula_bg.fill.solid()
    formula_bg.fill.fore_color.rgb = RGBColor(246, 248, 250)
    formula_bg.line.color.rgb = RGBColor(210, 218, 226)
    add_textbox(s, Inches(1.05), Inches(1.82), Inches(5.45), Inches(0.35), "Core intervention", 18, True, ACCENT)
    add_textbox(s, Inches(1.05), Inches(2.42), Inches(5.45), Inches(0.5), "l*  =  argmax_l AUC_val(l)", 21, False, TEXT)
    add_textbox(s, Inches(1.05), Inches(3.08), Inches(5.45), Inches(0.6), "h'_{l*,t}  =  h_{l*,t}  +  α · v_{l*}", 25, True, TEXT)
    add_textbox(s, Inches(1.05), Inches(3.92), Inches(5.45), Inches(0.4), "α ∈ {0, 0.8, 1.5, 2.5, 4}", 19, False, MUTED)

    notation_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.25), Inches(1.55), Inches(5.15), Inches(3.05))
    notation_bg.fill.solid()
    notation_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    notation_bg.line.color.rgb = RGBColor(210, 218, 226)
    add_textbox(s, Inches(7.55), Inches(1.82), Inches(4.55), Inches(0.35), "Notation", 18, True, ACCENT)
    add_textbox(
        s,
        Inches(7.55),
        Inches(2.42),
        Inches(4.55),
        Inches(1.75),
        "l*: validation-selected best layer\n"
        "h_{l*,t}: hidden state at layer l* and token t\n"
        "v_{l*}: direction learned at layer l*\n"
        "α: steering strength / coefficient\n"
        "h'_{l*,t}: modified hidden state used for generation",
        15,
        False,
        TEXT,
    )

    eval_bg = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.0), Inches(11.65), Inches(1.2))
    eval_bg.fill.solid()
    eval_bg.fill.fore_color.rgb = RGBColor(235, 243, 252)
    eval_bg.line.color.rgb = RGBColor(210, 218, 226)
    add_textbox(s, Inches(1.05), Inches(5.22), Inches(2.0), Inches(0.3), "Evaluation", 18, True, ACCENT)
    add_textbox(
        s,
        Inches(3.0),
        Inches(5.23),
        Inches(9.0),
        Inches(0.7),
        "Generate answers with α, parse stance labels, then compare\n"
        "P_α(stance | model, language, side) and Δ mean stance(α)",
        18,
        False,
        TEXT,
    )

    add_textbox(
        s,
        Inches(0.85),
        Inches(6.45),
        Inches(11.4),
        Inches(0.28),
        "Interpretation: only the best layer l* is intervened on; if stance shifts as α increases, that layer's direction is behaviorally steerable.",
        14,
        False,
        MUTED,
    )

    # 4
    s = blank(prs)
    add_title(s, "Experimental Scope", "The current deck uses the three instruct models only.")
    rows = [
        ("Models", "Qwen2.5-7B-Instruct; Mistral-7B-Instruct-v0.2; Llama-3-8B-Instruct"),
        ("Languages", "English, Italian, French, Spanish, German"),
        ("Interventions", "Left steering and right steering"),
        ("Coefficients", "0, 0.8, 1.5, 2.5, 4"),
        ("Items", "62 political compass-style statements per condition"),
    ]
    y = Inches(1.65)
    for key, val in rows:
        add_textbox(s, Inches(0.85), y, Inches(2.0), Inches(0.34), key, 16, True, ACCENT)
        add_textbox(s, Inches(3.0), y, Inches(9.3), Inches(0.34), val, 16)
        y += Inches(0.75)

    # 5
    s = blank(prs)
    add_title(s, "Method Overview", "A simple pipeline: answer, parse, detect, steer, compare.")
    add_numbered_cards(
        s,
        [
            ("Collect answers", "Ask each model to answer multilingual political statements under baseline and steered settings."),
            ("Parse stance", "Map outputs into four stance classes: strongly disagree, disagree, agree, strongly agree."),
            ("Compare shifts", "Measure how stance distributions and mean stance scores change as coefficient increases."),
        ],
    )

    # 6
    s = blank(prs)
    add_title(s, "Detection vs. Steering", "This distinction is central to the story.")
    add_bullets(
        s,
        Inches(0.75),
        Inches(1.75),
        Inches(11.7),
        Inches(4.6),
        [
            "Detection asks whether political response patterns are separable. A high AUC can mean the task creates a very clean ranking signal.",
            "Steering asks whether an intervention changes behavior. This is closer to a causal test.",
            "AUC can be near 1 in a controlled benchmark, but that does not automatically imply robust real-world political understanding.",
            "The current results suggest: detectable does not necessarily mean strongly steerable.",
        ],
        20,
    )

    # 7
    s = blank(prs)
    add_title(s, "Detection Result: Hidden-State Separability", "Best-layer AUC is high across models, but Qwen-left is the visibly weaker direction.")
    add_line(s, Inches(0.45), Inches(1.42), Inches(0.45), Inches(5.05), RGBColor(215, 220, 225), 1.3)
    add_story_detection_chart(s, Inches(0.65), Inches(1.55), Inches(6.2), Inches(3.05))
    add_line(s, Inches(7.05), Inches(1.42), Inches(7.05), Inches(5.05), RGBColor(215, 220, 225), 1.3)
    add_textbox(
        s,
        Inches(7.75),
        Inches(1.28),
        Inches(4.65),
        Inches(0.32),
        "Best layer used for steering",
        15,
        True,
        ACCENT,
    )
    add_story_detection_table(s, Inches(7.55), Inches(1.72), Inches(4.0), Inches(2.35))
    add_panel(
        s,
        Inches(7.55),
        Inches(4.35),
        Inches(4.55),
        Inches(1.05),
        "How to read this",
        "Best layer is the single layer selected by validation AUC and used for steering. The table reports the steering layer and its test AUC.",
        RGBColor(255, 242, 204),
    )
    add_panel(
        s,
        Inches(0.65),
        Inches(5.85),
        Inches(11.55),
        Inches(0.62),
        "Takeaway",
        "Detection is strong across all three models, but downstream steering still tests only the selected best-layer direction.",
        RGBColor(248, 249, 250),
    )
    add_textbox(
        s,
        Inches(0.75),
        Inches(6.68),
        Inches(11.6),
        Inches(0.28),
        "Source: story_detection_data.csv; Qwen values are parsed from the promptfixed steering logs so the AUC/layer matches the intervention.",
        9.5,
        False,
        MUTED,
    )

    # 8
    s = blank(prs)
    add_title(s, "Steering Intervention Design", "Two directions, five strengths, multiple languages.")
    add_bullets(
        s,
        Inches(0.85),
        Inches(1.7),
        Inches(5.7),
        Inches(4.6),
        [
            "Left steering: push activations in the learned left-leaning direction.",
            "Right steering: push activations in the opposite direction.",
            "Coefficient controls intervention strength.",
            "A true steering effect should produce a dose-response pattern rather than random fluctuation.",
        ],
        18,
    )
    add_textbox(s, Inches(7.1), Inches(2.0), Inches(4.8), Inches(0.45), "Expected evidence pattern", 20, True, ACCENT)
    add_bullets(
        s,
        Inches(7.1),
        Inches(2.65),
        Inches(4.8),
        Inches(2.6),
        [
            "coef increases -> stance score changes monotonically or near-monotonically",
            "left and right steering move in opposite directions",
            "effect repeats across at least some languages",
        ],
        17,
    )

    # 9
    s = blank(prs)
    add_title(s, "Output Parsing and Stance Labels", "Quality control comes before interpretation.")
    add_bullets(
        s,
        Inches(0.8),
        Inches(1.7),
        Inches(11.3),
        Inches(4.8),
        [
            "Each model answer is parsed into one of four ordered stance labels.",
            "Mean stance score summarizes movement, but it hides the underlying distribution.",
            "Stacked bars show the full distribution, including unparsed outputs.",
            "Radar plots show the shape of stance proportions at the strongest coefficient.",
        ],
        20,
    )

    # 10
    s = blank(prs)
    add_title(s, "Quality Control: Parse Rate", "Mistral requires caution because some language settings are less parseable.")
    add_image(s, FIG_DIR / "parse_rate_heatmap.png", Inches(0.7), Inches(1.45), Inches(11.9), Inches(5.45))

    # 11 – per-language steering tables
    lang_slide_info = [
        ("en", "English"),
        ("it", "Italian"),
        ("fr", "French"),
        ("es", "Spanish"),
        ("de", "German"),
    ]
    for lang_code, lang_label in lang_slide_info:
        s = blank(prs)
        add_title(
            s,
            f"Steering Results: {lang_label} Item and Coordinate Changes",
            f"Yellow alpha=4 rows are the {lang_label} rows repeated in the cross-lingual summary.",
        )
        add_image(
            s,
            STORY_FIG_DIR / f"fig2a_{lang_code}_item_coordinate_change_table.png",
            Inches(0.45),
            Inches(1.25),
            Inches(12.45),
            Inches(5.55),
        )
        add_textbox(
            s,
            Inches(0.75),
            Inches(6.93),
            Inches(11.8),
            Inches(0.28),
            "Reading: Changed items compare parsed stance labels with alpha=0; Delta Econ/Social and distance summarize movement in compass coordinates.",
            10,
            False,
            MUTED,
        )

    # 12
    s = blank(prs)
    add_title(
        s,
        "Cross-Lingual Item and Coordinate Changes",
        "At alpha=4, the strongest behavioral movement is model- and language-dependent.",
    )
    add_image(
        s,
        STORY_FIG_DIR / "fig2b_crosslingual_alpha4_item_coordinate_change_table.png",
        Inches(0.38),
        Inches(1.22),
        Inches(12.55),
        Inches(5.55),
    )
    add_textbox(
        s,
        Inches(0.75),
        Inches(6.93),
        Inches(11.8),
        Inches(0.28),
        "Takeaway: Llama-3 left steering changes many more items than Qwen or Mistral; right steering is weaker except for Llama-3 French.",
        10,
        False,
        MUTED,
    )

    # Stance distribution slides: 3 models x 5 languages
    model_slugs = {"Qwen": "qwen", "Mistral": "mistral", "Llama-3": "llama3"}
    for model_name, model_slug in model_slugs.items():
        for lang_code, lang_label in lang_slide_info:
            s = blank(prs)
            add_title(
                s,
                f"Stance Distribution: {model_name} {lang_label}",
                "Stacked bars show how answer-category proportions shift with increasing steering coefficient.",
            )
            add_image(
                s,
                STORY_FIG_DIR / f"figS1_focused_stance_distribution_{model_slug}_{lang_code}.png",
                Inches(0.4),
                Inches(1.15),
                Inches(12.5),
                Inches(5.7),
            )
            add_textbox(
                s,
                Inches(0.75),
                Inches(6.93),
                Inches(11.8),
                Inches(0.28),
                "This supplementary view shows answer-category proportions per steering coefficient; it is not a coordinate trajectory.",
                10,
                False,
                MUTED,
            )

    # Detailed views
    s = blank(prs)
    add_title(s, "Llama Detailed View", "The clearest case of behavioral movement under steering.")
    add_image(s, FIG_DIR / "stance_stacked_llama-3-8b-instruct.png", Inches(0.55), Inches(1.35), Inches(6.45), Inches(5.65))
    add_image(s, FIG_DIR / "radar_coef4_llama-3-8b-instruct.png", Inches(7.15), Inches(1.65), Inches(5.65), Inches(4.85))

    # 13
    s = blank(prs)
    add_title(s, "Qwen and Mistral Detailed View", "Two weaker/less direct steering stories.")
    add_image(s, FIG_DIR / "radar_coef4_qwen2_5-7b-instruct.png", Inches(0.7), Inches(1.75), Inches(5.75), Inches(4.9))
    add_image(s, FIG_DIR / "radar_coef4_mistral-7b-instruct-v0_2.png", Inches(6.85), Inches(1.75), Inches(5.75), Inches(4.9))

    # 14
    s = blank(prs)
    add_title(s, "Current Interpretation and Next Steps")
    add_bullets(
        s,
        Inches(0.75),
        Inches(1.55),
        Inches(11.8),
        Inches(4.9),
        [
            "Main story: political stance behavior is model-dependent, language-dependent, and intervention-dependent.",
            "Key claim to sharpen: high detectability does not guarantee high steerability.",
            "Current result table separates item-level answer changes from coordinate-level compass movement.",
            "Next writing step: turn each slide into a paragraph-level explanation for the paper/report.",
        ],
        20,
    )

    prs.save(OUT_PATH)
    print(OUT_PATH)


if __name__ == "__main__":
    build()
